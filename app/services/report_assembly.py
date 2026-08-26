"""Complete PDF report assembly: selects the latest authoritative Event
Report document, converts DOCX/PPTX to PDF (preserving paragraphs, tables,
embedded images, and authorship) or keeps an already-PDF source unchanged,
appends testimonial/programme documents, and falls back to a generated
summary when no authoritative report exists. See PLAN.md "Complete ICC
report generation"."""

from __future__ import annotations

import hashlib
import io
import tempfile
from datetime import datetime, timezone

import docx
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image as RLImage
from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from app.database import db
from app.models.erp import DocumentRecord, ReportSnapshot, SessionAttendance
from app.services.authorization import has_permission
from app.services.documents import compute_availability
from app.services.drive import download_drive_file

APPENDIX_CATEGORIES = ["Testimonial", "Program Details", "Inauguration Schedule", "Valedictory Schedule"]
SUPPORTED_EXTENSIONS = {"pdf", "docx", "pptx", "jpg", "jpeg", "png"}
MAX_FETCH_BYTES = 100 * 1024 * 1024
FETCH_TIMEOUT_SECONDS = 30


class ReportIncompleteError(Exception):
    def __init__(self, preflight):
        self.preflight = preflight
        super().__init__("Report dependencies are incomplete; acknowledge with allow_incomplete to continue.")


def _extension(document: DocumentRecord) -> str:
    name = document.drive_name or document.title or ""
    return name.rsplit(".", 1)[-1].lower() if "." in name else ""


def select_authoritative_report(project) -> DocumentRecord | None:
    candidates = (
        DocumentRecord.query.filter_by(project_id=project.id, category="Event Report")
        .order_by(DocumentRecord.drive_modified_at.desc(), DocumentRecord.created_at.desc())
        .all()
    )
    return candidates[0] if candidates else None


def _appendix_documents(project) -> list[DocumentRecord]:
    return (
        DocumentRecord.query.filter(
            DocumentRecord.project_id == project.id,
            DocumentRecord.category.in_(APPENDIX_CATEGORIES),
        ).order_by(DocumentRecord.category, DocumentRecord.created_at).all()
    )


def _dependency_status(document: DocumentRecord) -> str:
    availability = compute_availability(document)
    if availability == "Missing":
        return "missing"
    if availability == "Inaccessible":
        return "inaccessible"
    if _extension(document) not in SUPPORTED_EXTENSIONS:
        return "unsupported"
    return "included"


def preflight_report(project) -> dict:
    """Report each report dependency as included/missing/inaccessible/unsupported."""
    dependencies = []
    authoritative = select_authoritative_report(project)
    if authoritative is not None:
        dependencies.append({
            "role": "authoritative_report", "document_public_id": authoritative.public_id,
            "title": authoritative.title, "status": _dependency_status(authoritative),
        })
    for document in _appendix_documents(project):
        dependencies.append({
            "role": "appendix", "document_public_id": document.public_id,
            "title": document.title, "status": _dependency_status(document),
        })
    buckets = {status: [] for status in ("included", "missing", "inaccessible", "unsupported")}
    for dependency in dependencies:
        buckets[dependency["status"]].append(dependency)
    return {
        "dependencies": dependencies,
        "included": buckets["included"],
        "missing": buckets["missing"],
        "inaccessible": buckets["inaccessible"],
        "unsupported": buckets["unsupported"],
        "has_authoritative_report": authoritative is not None,
        "complete": not (buckets["missing"] or buckets["inaccessible"] or buckets["unsupported"]),
    }


def _fetch_document_bytes(document: DocumentRecord) -> bytes:
    """Fetch a Drive binary into a temporary file for the duration of
    processing only; the file is deleted as soon as the `with` block exits.
    Enforces a size ceiling so a single dependency cannot exhaust memory."""
    export_mime = None
    if document.drive_mime_type and document.drive_mime_type.startswith("application/vnd.google-apps"):
        export_mime = "application/pdf"
    content = download_drive_file(document.drive_file_id, export_mime_type=export_mime)
    if len(content) > MAX_FETCH_BYTES:
        raise ValueError(f"{document.title}: fetched content exceeds the report-assembly size limit.")
    with tempfile.NamedTemporaryFile(delete=True) as handle:
        handle.write(content)
        handle.flush()
        handle.seek(0)
        return handle.read()


def _styles():
    return getSampleStyleSheet()


def _docx_flowables(content: bytes, *, seen_image_hashes: set) -> list:
    document = docx.Document(io.BytesIO(content))
    styles = _styles()
    flowables = []
    author = (document.core_properties.author or "").strip()
    if author:
        flowables.append(Paragraph(f"<i>Prepared by: {author}</i>", styles["Normal"]))
        flowables.append(Spacer(1, 12))
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = "Heading2" if (paragraph.style and paragraph.style.name or "").startswith("Heading") else "Normal"
        flowables.append(Paragraph(text, styles[style_name]))
        flowables.append(Spacer(1, 6))
    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if rows:
            rl_table = Table(rows)
            rl_table.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 0.5, "grey"), ("FONTSIZE", (0, 0), (-1, -1), 8)]))
            flowables.append(rl_table)
            flowables.append(Spacer(1, 12))
    for rel in document.part.rels.values():
        if "image" not in rel.reltype:
            continue
        blob = rel.target_part.blob
        digest = hashlib.sha256(blob).hexdigest()
        if digest in seen_image_hashes:
            continue
        seen_image_hashes.add(digest)
        try:
            image = RLImage(io.BytesIO(blob), width=5 * inch, height=3.5 * inch, kind="proportional")
        except Exception:
            continue
        flowables.append(image)
        flowables.append(Spacer(1, 12))
    return flowables


def _pptx_flowables(content: bytes, *, seen_image_hashes: set) -> list:
    presentation = Presentation(io.BytesIO(content))
    styles = _styles()
    flowables = []
    for index, slide in enumerate(presentation.slides, start=1):
        if index > 1:
            flowables.append(PageBreak())
        flowables.append(Paragraph(f"Slide {index}", styles["Heading3"]))
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                flowables.append(Paragraph(shape.text_frame.text.strip(), styles["Normal"]))
                flowables.append(Spacer(1, 6))
            if shape.shape_type == 13 and hasattr(shape, "image"):  # PICTURE
                blob = shape.image.blob
                digest = hashlib.sha256(blob).hexdigest()
                if digest in seen_image_hashes:
                    continue
                seen_image_hashes.add(digest)
                try:
                    flowables.append(RLImage(io.BytesIO(blob), width=5 * inch, height=3.5 * inch, kind="proportional"))
                    flowables.append(Spacer(1, 12))
                except Exception:
                    continue
    return flowables


def _image_flowables(content: bytes) -> list:
    return [RLImage(io.BytesIO(content), width=6 * inch, height=8 * inch, kind="proportional")]


def _render_section_pdf(document: DocumentRecord, content: bytes, *, seen_image_hashes: set) -> bytes:
    extension = _extension(document)
    if extension == "pdf":
        return content
    buffer = io.BytesIO()
    doc_template = SimpleDocTemplate(buffer, pagesize=LETTER)
    styles = _styles()
    flowables = [Paragraph(document.title, styles["Title"]), Spacer(1, 12)]
    if extension == "docx":
        flowables.extend(_docx_flowables(content, seen_image_hashes=seen_image_hashes))
    elif extension == "pptx":
        flowables.extend(_pptx_flowables(content, seen_image_hashes=seen_image_hashes))
    elif extension in {"jpg", "jpeg", "png"}:
        flowables.extend(_image_flowables(content))
    else:
        raise ValueError(f"Unsupported document type for report assembly: {extension}")
    doc_template.build(flowables)
    buffer.seek(0)
    return buffer.read()


def _merge_pdf_sections(sections: list[bytes]) -> bytes:
    writer = PdfWriter()
    for section in sections:
        reader = PdfReader(io.BytesIO(section))
        for page in reader.pages:
            writer.add_page(page)
    output = io.BytesIO()
    writer.write(output)
    output.seek(0)
    return output.read()


def _fallback_summary_pdf(project) -> bytes:
    buffer = io.BytesIO()
    doc_template = SimpleDocTemplate(buffer, pagesize=LETTER)
    styles = _styles()
    session_ids = [session.id for session in project.sessions if session.is_active]
    present_count = (
        SessionAttendance.query.filter(SessionAttendance.session_id.in_(session_ids), SessionAttendance.status == "Present").count()
        if session_ids else 0
    )
    flowables = [
        Paragraph(project.title, styles["Title"]),
        Spacer(1, 12),
        Paragraph(f"No authoritative Event Report document is available. This is a generated summary.", styles["Italic"]),
        Spacer(1, 12),
        Paragraph(f"Category: {project.category}", styles["Normal"]),
        Paragraph(f"Dates: {project.start_date} to {project.end_date}", styles["Normal"]),
        Paragraph(f"Venue: {project.venue or '—'}", styles["Normal"]),
        Paragraph(f"Target audience: {project.target_audience or '—'}", styles["Normal"]),
        Paragraph(f"Expected reach: {project.expected_reach if project.expected_reach is not None else '—'}", styles["Normal"]),
        Paragraph(f"Actual reach: {project.actual_reach if project.actual_reach is not None else '—'}", styles["Normal"]),
        Paragraph(f"Recorded present attendance: {present_count}", styles["Normal"]),
        Spacer(1, 12),
        Paragraph("Schedule", styles["Heading2"]),
    ]
    for session in sorted(project.sessions, key=lambda s: s.starts_at):
        if not session.is_active:
            continue
        when = "All day" if session.is_all_day else f"{session.starts_at.strftime('%I:%M %p')}–{session.ends_at.strftime('%I:%M %p')}"
        flowables.append(Paragraph(f"{session.starts_at.strftime('%b %d, %Y')} — {session.title} ({when})", styles["Normal"]))
    doc_template.build(flowables)
    buffer.seek(0)
    return buffer.read()


def assemble_complete_report(project, actor, *, allow_incomplete=False, fetch_override=None):
    """Return (pdf_bytes, ReportSnapshot). Raises ReportIncompleteError if
    any dependency is missing/inaccessible/unsupported and
    `allow_incomplete` is not set -- the caller (route) is responsible for
    surfacing the warning and requiring an explicit re-request."""
    preflight = preflight_report(project)
    if not preflight["complete"] and not allow_incomplete:
        raise ReportIncompleteError(preflight)

    fetch = fetch_override or _fetch_document_bytes
    seen_image_hashes: set = set()
    sections: list[bytes] = []
    included_documents = []

    authoritative = select_authoritative_report(project)
    if authoritative is not None and _dependency_status(authoritative) == "included":
        content = fetch(authoritative)
        sections.append(_render_section_pdf(authoritative, content, seen_image_hashes=seen_image_hashes))
        included_documents.append(_source_reference(authoritative))
    else:
        sections.append(_fallback_summary_pdf(project))

    for document in _appendix_documents(project):
        if _dependency_status(document) != "included":
            continue
        if document.permission_classification == "Restricted" and not has_permission(actor, "sensitive_links", project, sensitive=True):
            continue  # never embed a restricted document without the existing sensitive-link permission
        content = fetch(document)
        sections.append(_render_section_pdf(document, content, seen_image_hashes=seen_image_hashes))
        included_documents.append(_source_reference(document))

    pdf_bytes = _merge_pdf_sections(sections)
    snapshot = _persist_snapshot(project, actor, included_documents, preflight, allow_incomplete)
    return pdf_bytes, snapshot


def _source_reference(document: DocumentRecord) -> dict:
    return {
        "document_public_id": document.public_id,
        "title": document.title,
        "category": document.category,
        "drive_modified_at": document.drive_modified_at.isoformat() if document.drive_modified_at else None,
        "checksum_sha256": document.checksum_sha256,
    }


def _persist_snapshot(project, actor, included_documents, preflight, allow_incomplete) -> ReportSnapshot:
    version = (
        db.session.query(db.func.max(ReportSnapshot.version))
        .filter_by(project_id=project.id, report_type="complete_pdf")
        .scalar() or 0
    ) + 1
    snapshot = ReportSnapshot(
        project_id=project.id, report_type="complete_pdf", title=f"{project.title} — Complete Report",
        version=version,
        snapshot_json={
            "included_documents": included_documents,
            "inferred_metrics": {
                "actual_reach": project.actual_reach,
                "expected_reach": project.expected_reach,
            },
            "acknowledged_incomplete": allow_incomplete and not preflight["complete"],
            "preflight_at_generation": preflight,
            "generated_at": datetime.now(timezone.utc).isoformat(),
        },
        source_references=[doc["document_public_id"] for doc in included_documents],
        generated_by_id=getattr(actor, "id", None),
    )
    db.session.add(snapshot)
    db.session.commit()
    return snapshot
